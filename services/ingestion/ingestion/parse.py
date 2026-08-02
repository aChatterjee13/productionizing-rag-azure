"""Format-aware parsing into ordered :class:`~ragcore.models.document.ParsedBlock`s.

Every parser produces the same shape — a flat list of blocks in reading order, each
carrying its structural ``kind``, its heading ``level`` (headings only), the
``section_path`` breadcrumb in effect, and a ``page`` number where the format is
paginated. That is what lets one structure-aware chunker in :mod:`ingestion.chunk`
serve PDF, DOCX, PPTX, XLSX, HTML, Markdown, CSV and plain text without special cases.

Supported inputs and the library each uses:

| format | library | notes |
|---|---|---|
| PDF | ``pdfplumber``, else ``pypdf`` | per-page blocks, tables extracted |
| DOCX | ``python-docx`` | ``Heading N`` styles become heading levels |
| PPTX | ``python-pptx`` | one slide = one section, slide number = page |
| XLSX | ``openpyxl`` | one sheet = one section, rows batched into tables |
| HTML | ``selectolax``, else ``beautifulsoup4`` | script/style stripped |
| Markdown | built in | ATX/setext headings, fences, lists, pipe tables |
| CSV/TSV | built in | header repeated on every table block |
| text | built in | blank-line separated paragraphs |

Tables are never flattened into prose: a table becomes one ``TABLE`` block rendered as
a pipe table, which the chunker keeps intact so a row never loses its header.
"""

from __future__ import annotations

import csv
import io
import re
from collections.abc import Iterable, Sequence
from typing import Any

from ingestion.chunk import estimate_tokens
from ragcore.errors import RagError
from ragcore.logging import get_logger
from ragcore.models.document import (
    BlockKind,
    ParsedBlock,
    ParsedDocument,
    SourceDocument,
)
from ragcore.settings import Settings, get_settings

__all__ = [
    "ParseError",
    "SectionTracker",
    "normalise_text",
    "parse_csv",
    "parse_document",
    "parse_docx",
    "parse_html",
    "parse_markdown",
    "parse_pdf",
    "parse_pptx",
    "parse_text",
    "parse_xlsx",
    "render_table",
]

_log = get_logger(__name__)

#: Markdown ATX heading, e.g. ``### Overview``.
_ATX_RE = re.compile(r"^(#{1,6})\s+(.*?)\s*#*\s*$")

#: Ordered or unordered list item.
_LIST_RE = re.compile(r"^\s{0,3}(?:[-*+]\s+|\d{1,3}[.)]\s+)(.*)$")

#: Numbered section heading in extracted PDF text, e.g. ``4.2.1 Retention``.
_NUMBERED_HEADING_RE = re.compile(r"^(\d+(?:\.\d+){0,4})[.)]?\s+(\S.{0,110})$")

#: Collapse runs of whitespace while preserving paragraph breaks.
_WS_RE = re.compile(r"[ \t\r\f\v]+")

#: Maximum characters in a line still eligible to be a heuristic heading.
_HEADING_MAX_CHARS = 90


class ParseError(RagError):
    """A document could not be parsed into blocks.

    Raised when the payload is empty or a required parsing library is not installed.
    The pipeline records the document as failed and moves on.
    """

    status_code = 422
    code = "parse_error"


class SectionTracker:
    """Maintain the heading breadcrumb while walking a document.

    A heading at level ``L`` pops every deeper-or-equal entry off the stack and pushes
    itself, so ``section_path`` always reads root-first and never contains a stale
    sibling.
    """

    def __init__(self) -> None:
        """Start with an empty breadcrumb."""
        self._stack: list[tuple[int, str]] = []

    def push(self, level: int, title: str) -> list[str]:
        """Record a heading and return the breadcrumb it establishes.

        Args:
            level: Heading depth, 1 = top level.
            title: Heading text.

        Returns:
            The section path including this heading.
        """
        depth = max(1, level)
        self._stack = [entry for entry in self._stack if entry[0] < depth]
        cleaned = title.strip()
        if cleaned:
            self._stack.append((depth, cleaned))
        return self.current()

    def current(self) -> list[str]:
        """Breadcrumb in effect for the next content block.

        Returns:
            A copy of the current section path.
        """
        return [title for _, title in self._stack]

    def reset(self) -> None:
        """Clear the breadcrumb, e.g. at a new page or sheet."""
        self._stack = []


class _BlockBuilder:
    """Accumulate blocks with monotonic ``order`` and the live section path."""

    def __init__(self) -> None:
        """Start an empty block list with a fresh section tracker."""
        self.blocks: list[ParsedBlock] = []
        self.sections = SectionTracker()

    def heading(self, text: str, level: int, *, page: int | None = None) -> None:
        """Append a heading block and advance the breadcrumb.

        Args:
            text: Heading text.
            level: Heading depth, 1 = top level.
            page: Page number for paginated formats.
        """
        cleaned = normalise_text(text)
        if not cleaned:
            return
        path = self.sections.push(level, cleaned)
        self._append(BlockKind.HEADING, cleaned, page=page, level=level, path=path)

    def block(
        self,
        kind: BlockKind,
        text: str,
        *,
        page: int | None = None,
        metadata: dict[str, Any] | None = None,
    ) -> None:
        """Append a content block under the current breadcrumb.

        Args:
            kind: Structural role.
            text: Block text.
            page: Page number for paginated formats.
            metadata: Parser extras such as table shape or code language.
        """
        cleaned = normalise_text(text)
        if not cleaned:
            return
        self._append(
            kind, cleaned, page=page, path=self.sections.current(), metadata=metadata
        )

    def _append(
        self,
        kind: BlockKind,
        text: str,
        *,
        page: int | None,
        path: list[str],
        level: int | None = None,
        metadata: dict[str, Any] | None = None,
    ) -> None:
        """Append a fully-formed block.

        Args:
            kind: Structural role.
            text: Normalised text.
            page: Page number, when paginated.
            path: Section breadcrumb.
            level: Heading depth for heading blocks.
            metadata: Parser extras.
        """
        self.blocks.append(
            ParsedBlock(
                kind=kind,
                text=text,
                order=len(self.blocks),
                level=level,
                page=page,
                section_path=list(path),
                metadata=dict(metadata or {}),
            )
        )


def normalise_text(text: str) -> str:
    """Collapse horizontal whitespace and trim, preserving newlines.

    Args:
        text: Raw extracted text.

    Returns:
        The normalised text, or "" when nothing is left.
    """
    if not text:
        return ""
    unwrapped = text.replace("\u00a0", " ")
    lines = [_WS_RE.sub(" ", line).strip() for line in unwrapped.split("\n")]
    return "\n".join(line for line in lines if line != "").strip()


def render_table(rows: Sequence[Sequence[Any]]) -> str:
    """Render tabular data as a Markdown pipe table.

    Args:
        rows: Rows of cells; the first row is treated as the header.

    Returns:
        The rendered table, or "" when there is nothing to render.
    """
    cleaned = [
        ["" if cell is None else _WS_RE.sub(" ", str(cell)).strip() for cell in row]
        for row in rows
        if any(cell is not None and str(cell).strip() for cell in row)
    ]
    if not cleaned:
        return ""
    width = max(len(row) for row in cleaned)
    padded = [[*row, *([""] * (width - len(row)))] for row in cleaned]
    header, *body = padded
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


# --------------------------------------------------------------------- dispatch
def parse_document(
    doc: SourceDocument, settings: Settings | None = None
) -> ParsedDocument:
    """Parse a fetched document into ordered blocks.

    Args:
        doc: The fetched document, carrying bytes or text plus its media type.
        settings: Process settings; ``get_settings()`` is used when omitted.

    Returns:
        A :class:`ParsedDocument` whose ``blocks`` are ordered, breadcrumbed and
        page-numbered, and whose ``title``/``author``/``page_count`` are filled in
        from the format's own metadata when the source did not supply them.

    Raises:
        ParseError: If the document carries no payload, or the library needed for its
            format is not installed.
    """
    active = settings or get_settings()
    if not doc.has_content:
        msg = "document has no content to parse"
        raise ParseError(msg, detail={"document_id": doc.document_id})

    media_type = (doc.media_type or "").split(";")[0].strip().lower()
    name = (doc.filename or doc.source_uri or "").lower()
    extra: dict[str, Any] = {}

    if media_type == "application/pdf" or name.endswith(".pdf"):
        builder, extra = parse_pdf(_require_bytes(doc), active)
    elif "wordprocessingml" in media_type or name.endswith(".docx"):
        builder, extra = parse_docx(_require_bytes(doc))
    elif "presentationml" in media_type or name.endswith(".pptx"):
        builder, extra = parse_pptx(_require_bytes(doc))
    elif "spreadsheetml" in media_type or name.endswith((".xlsx", ".xlsm")):
        builder, extra = parse_xlsx(_require_bytes(doc), active)
    elif media_type in {"text/html", "application/xhtml+xml"} or name.endswith(
        (".html", ".htm")
    ):
        builder, extra = parse_html(doc.text_or_empty())
    elif media_type in {"text/markdown"} or name.endswith((".md", ".markdown")):
        builder = parse_markdown(doc.text_or_empty())
    elif media_type in {"text/csv", "text/tab-separated-values"} or name.endswith(
        (".csv", ".tsv")
    ):
        builder = parse_csv(doc.text_or_empty(), active, tab=name.endswith(".tsv"))
    else:
        builder = parse_text(doc.text_or_empty())

    blocks = builder.blocks
    if not blocks:
        msg = "parser produced no blocks"
        raise ParseError(msg, detail={"document_id": doc.document_id})

    title = doc.title.strip() or str(extra.get("title") or "").strip()
    if not title:
        title = _first_heading(blocks) or _title_from_name(doc.filename)

    return ParsedDocument(
        document_id=doc.document_id,
        tenant_id=doc.tenant_id,
        source_id=doc.source_id,
        source_type=doc.source_type,
        source_uri=doc.source_uri,
        title=title,
        blocks=blocks,
        doc_type=doc.doc_type,
        language=doc.language or "en",
        author=doc.author or (str(extra["author"]) if extra.get("author") else None),
        tags=list(doc.tags),
        page_count=extra.get("page_count"),
        access_control=doc.access_control,
        content_sha256=doc.content_sha256,
        source_modified_at=doc.source_modified_at,
        effective_from=doc.effective_from,
        effective_to=doc.effective_to,
        metadata={**doc.metadata, "media_type": media_type},
    )


def _require_bytes(doc: SourceDocument) -> bytes:
    """Return a binary document's payload.

    Args:
        doc: The fetched document.

    Returns:
        The raw bytes.

    Raises:
        ParseError: If the document only carries text, which a binary parser cannot
            use.
    """
    if doc.content_bytes is None:
        msg = "binary format requires content_bytes"
        raise ParseError(msg, detail={"document_id": doc.document_id})
    return doc.content_bytes


def _first_heading(blocks: Iterable[ParsedBlock]) -> str:
    """Find the first heading in a block list.

    Args:
        blocks: Parsed blocks.

    Returns:
        The heading text, or "" when the document has no headings.
    """
    for block in blocks:
        if block.kind is BlockKind.HEADING:
            return block.text
    return ""


def _title_from_name(filename: str | None) -> str:
    """Derive a readable title from a file name.

    Args:
        filename: Original file name.

    Returns:
        The stem with separators replaced by spaces, or ``"Untitled"``.
    """
    if not filename:
        return "Untitled"
    stem = filename.rsplit("/", 1)[-1].rsplit(".", 1)[0]
    cleaned = stem.replace("_", " ").replace("-", " ").strip()
    return cleaned or "Untitled"


# ------------------------------------------------------------------- plain text
def parse_text(text: str) -> _BlockBuilder:
    """Parse plain text into paragraph blocks.

    Args:
        text: Decoded document text.

    Returns:
        A builder holding one block per blank-line-separated paragraph.
    """
    builder = _BlockBuilder()
    for paragraph in re.split(r"\n\s*\n", text):
        builder.block(BlockKind.PARAGRAPH, paragraph)
    return builder


# --------------------------------------------------------------------- markdown
def parse_markdown(text: str) -> _BlockBuilder:
    """Parse Markdown into headings, lists, code fences, tables and paragraphs.

    Args:
        text: Markdown source.

    Returns:
        A builder holding the parsed blocks.
    """
    builder = _BlockBuilder()
    lines = text.replace("\r\n", "\n").split("\n")
    buffer: list[str] = []
    table: list[str] = []
    fence: str | None = None
    fence_lines: list[str] = []
    fence_language = ""

    def flush_paragraph() -> None:
        if buffer:
            builder.block(BlockKind.PARAGRAPH, "\n".join(buffer))
            buffer.clear()

    def flush_table() -> None:
        if table:
            builder.block(
                BlockKind.TABLE,
                "\n".join(table),
                metadata={"rows": len(table)},
            )
            table.clear()

    index = 0
    while index < len(lines):
        line = lines[index]
        stripped = line.strip()

        if fence is not None:
            if stripped.startswith(fence):
                builder.block(
                    BlockKind.CODE,
                    "\n".join(fence_lines),
                    metadata={"language": fence_language},
                )
                fence = None
                fence_lines = []
                fence_language = ""
            else:
                fence_lines.append(line)
            index += 1
            continue

        if stripped.startswith(("```", "~~~")):
            flush_paragraph()
            flush_table()
            fence = stripped[:3]
            fence_language = stripped[3:].strip()
            index += 1
            continue

        atx = _ATX_RE.match(stripped)
        if atx:
            flush_paragraph()
            flush_table()
            builder.heading(atx.group(2), len(atx.group(1)))
            index += 1
            continue

        # Setext heading: text underlined with === or ---.
        if (
            stripped
            and index + 1 < len(lines)
            and re.fullmatch(r"=+|-{2,}", lines[index + 1].strip())
        ):
            flush_paragraph()
            flush_table()
            level = 1 if lines[index + 1].strip().startswith("=") else 2
            builder.heading(stripped, level)
            index += 2
            continue

        if stripped.startswith("|") and stripped.endswith("|"):
            flush_paragraph()
            table.append(stripped)
            index += 1
            continue
        flush_table()

        if stripped.startswith(">"):
            flush_paragraph()
            builder.block(BlockKind.QUOTE, stripped.lstrip("> "))
            index += 1
            continue

        list_match = _LIST_RE.match(line)
        if list_match:
            flush_paragraph()
            builder.block(BlockKind.LIST_ITEM, list_match.group(1))
            index += 1
            continue

        if not stripped:
            flush_paragraph()
        else:
            buffer.append(stripped)
        index += 1

    if fence is not None and fence_lines:
        builder.block(
            BlockKind.CODE,
            "\n".join(fence_lines),
            metadata={"language": fence_language},
        )
    flush_paragraph()
    flush_table()
    return builder


# --------------------------------------------------------------------------- csv
def parse_csv(text: str, settings: Settings, *, tab: bool = False) -> _BlockBuilder:
    """Parse delimited text into table blocks with a repeated header.

    Args:
        text: CSV or TSV content.
        settings: Process settings; ``chunk_max_tokens`` bounds one table block.
        tab: True for tab-separated input.

    Returns:
        A builder holding one or more table blocks.
    """
    builder = _BlockBuilder()
    delimiter = "\t" if tab else ","
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows = [row for row in reader if any(cell.strip() for cell in row)]
    if not rows:
        return builder
    header, *body = rows
    for batch in _batch_rows(body, header, settings):
        builder.block(
            BlockKind.TABLE,
            render_table([header, *batch]),
            metadata={"rows": len(batch), "columns": len(header)},
        )
    if not body:
        builder.block(
            BlockKind.TABLE,
            render_table([header]),
            metadata={"rows": 0, "columns": len(header)},
        )
    return builder


def _batch_rows(
    body: Sequence[Sequence[Any]], header: Sequence[Any], settings: Settings
) -> list[list[Sequence[Any]]]:
    """Group table rows so no single table block exceeds the chunk ceiling.

    Args:
        body: Data rows.
        header: Header row, repeated in every batch and therefore counted once per
            batch.
        settings: Process settings; ``chunk_max_tokens`` is the ceiling.

    Returns:
        A list of row batches.
    """
    header_cost = estimate_tokens(render_table([header]))
    ceiling = max(1, settings.chunk_max_tokens - header_cost)
    batches: list[list[Sequence[Any]]] = []
    current: list[Sequence[Any]] = []
    used = 0
    for row in body:
        rendered = " ".join("" if cell is None else str(cell) for cell in row)
        cost = estimate_tokens(rendered)
        if current and used + cost > ceiling:
            batches.append(current)
            current = []
            used = 0
        current.append(row)
        used += cost
    if current:
        batches.append(current)
    return batches


# -------------------------------------------------------------------------- html
def parse_html(html: str) -> tuple[_BlockBuilder, dict[str, Any]]:
    """Parse HTML into structural blocks.

    Args:
        html: Page markup.

    Returns:
        A ``(builder, extra)`` pair; ``extra`` may carry ``title``.
    """
    try:
        from selectolax.parser import HTMLParser
    except ImportError:
        return _parse_html_bs4(html)
    tree = HTMLParser(html)
    for tag in ("script", "style", "noscript", "template", "svg"):
        for node in tree.css(tag):
            node.decompose()
    builder = _BlockBuilder()
    title_node = tree.css_first("title")
    extra: dict[str, Any] = {}
    if title_node is not None and title_node.text():
        extra["title"] = normalise_text(title_node.text())

    body = tree.css_first("body") or tree.root
    if body is None:
        return builder, extra
    for node in body.iter(include_text=False):
        _emit_html_node(builder, node.tag, node)
    if not builder.blocks:
        builder.block(BlockKind.PARAGRAPH, tree.text(separator="\n"))
    return builder, extra


def _emit_html_node(builder: _BlockBuilder, tag: str, node: Any) -> None:
    """Emit blocks for one top-level HTML node.

    Args:
        builder: Block builder.
        tag: Lower-case tag name.
        node: The selectolax node.
    """
    if tag in {"h1", "h2", "h3", "h4", "h5", "h6"}:
        builder.heading(node.text(), int(tag[1]))
        return
    if tag == "table":
        rows = [
            [cell.text(strip=True) for cell in row.css("th, td")]
            for row in node.css("tr")
        ]
        builder.block(BlockKind.TABLE, render_table(rows), metadata={"rows": len(rows)})
        return
    if tag in {"ul", "ol"}:
        for item in node.css("li"):
            builder.block(BlockKind.LIST_ITEM, item.text())
        return
    if tag in {"pre", "code"}:
        builder.block(BlockKind.CODE, node.text())
        return
    if tag == "blockquote":
        builder.block(BlockKind.QUOTE, node.text())
        return
    if tag in {"figcaption", "caption"}:
        builder.block(BlockKind.CAPTION, node.text())
        return
    # Only emit leaf-ish containers so text is not duplicated by an ancestor.
    if tag in {"p", "article", "section", "div", "main", "td"} and not node.css(
        "p, h1, h2, h3, h4, h5, h6, table, ul, ol, pre, blockquote"
    ):
        builder.block(BlockKind.PARAGRAPH, node.text())


def _parse_html_bs4(html: str) -> tuple[_BlockBuilder, dict[str, Any]]:
    """Parse HTML with BeautifulSoup, or strip tags when neither library exists.

    Args:
        html: Page markup.

    Returns:
        A ``(builder, extra)`` pair.
    """
    builder = _BlockBuilder()
    extra: dict[str, Any] = {}
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        _log.warning("parse.html_no_library")
        stripped = re.sub(r"<[^>]+>", " ", html)
        builder.block(BlockKind.PARAGRAPH, stripped)
        return builder, extra

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript", "template", "svg"]):
        tag.decompose()
    if soup.title and soup.title.string:
        extra["title"] = normalise_text(str(soup.title.string))
    root = soup.body or soup
    for element in root.find_all(
        ["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "table", "pre", "blockquote"]
    ):
        name = element.name
        if name.startswith("h") and len(name) == 2 and name[1].isdigit():
            builder.heading(element.get_text(" ", strip=True), int(name[1]))
        elif name == "table":
            rows = [
                [cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])]
                for row in element.find_all("tr")
            ]
            builder.block(
                BlockKind.TABLE, render_table(rows), metadata={"rows": len(rows)}
            )
        elif name == "li":
            builder.block(BlockKind.LIST_ITEM, element.get_text(" ", strip=True))
        elif name == "pre":
            builder.block(BlockKind.CODE, element.get_text("\n", strip=True))
        elif name == "blockquote":
            builder.block(BlockKind.QUOTE, element.get_text(" ", strip=True))
        else:
            builder.block(BlockKind.PARAGRAPH, element.get_text(" ", strip=True))
    return builder, extra


# --------------------------------------------------------------------------- pdf
def parse_pdf(
    payload: bytes, settings: Settings
) -> tuple[_BlockBuilder, dict[str, Any]]:
    """Parse a PDF into per-page blocks, extracting tables where possible.

    ``pdfplumber`` is preferred because it recovers table structure; ``pypdf`` is the
    fallback and yields text only.

    Args:
        payload: Raw PDF bytes.
        settings: Process settings, used to bound table blocks.

    Returns:
        A ``(builder, extra)`` pair; ``extra`` carries ``page_count``, and ``title``
        and ``author`` when the PDF declares them.

    Raises:
        ParseError: If neither pdfplumber nor pypdf is installed.
    """
    try:
        import pdfplumber
    except ImportError:
        return _parse_pdf_pypdf(payload)

    builder = _BlockBuilder()
    extra: dict[str, Any] = {}
    with pdfplumber.open(io.BytesIO(payload)) as pdf:
        metadata = pdf.metadata or {}
        if metadata.get("Title"):
            extra["title"] = normalise_text(str(metadata["Title"]))
        if metadata.get("Author"):
            extra["author"] = normalise_text(str(metadata["Author"]))
        extra["page_count"] = len(pdf.pages)
        for index, page in enumerate(pdf.pages, start=1):
            for table in page.extract_tables() or []:
                if not table:
                    continue
                header, *body = table
                for batch in _batch_rows(body, header, settings):
                    rendered = render_table([header, *batch])
                    if rendered:
                        builder.block(
                            BlockKind.TABLE,
                            rendered,
                            page=index,
                            metadata={"rows": len(batch)},
                        )
            _emit_pdf_text(builder, page.extract_text() or "", index)
    return builder, extra


def _parse_pdf_pypdf(payload: bytes) -> tuple[_BlockBuilder, dict[str, Any]]:
    """Parse a PDF with pypdf when pdfplumber is unavailable.

    Args:
        payload: Raw PDF bytes.

    Returns:
        A ``(builder, extra)`` pair.

    Raises:
        ParseError: If pypdf is not installed either.
    """
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - optional install
        msg = "PDF parsing needs pdfplumber or pypdf"
        raise ParseError(msg) from exc

    builder = _BlockBuilder()
    reader = PdfReader(io.BytesIO(payload))
    extra: dict[str, Any] = {"page_count": len(reader.pages)}
    info = reader.metadata
    if info is not None:
        if info.title:
            extra["title"] = normalise_text(str(info.title))
        if info.author:
            extra["author"] = normalise_text(str(info.author))
    for index, page in enumerate(reader.pages, start=1):
        _emit_pdf_text(builder, page.extract_text() or "", index)
    return builder, extra


def _emit_pdf_text(builder: _BlockBuilder, text: str, page: int) -> None:
    """Turn one page of extracted PDF text into blocks.

    PDFs carry no structural markup, so headings are inferred: a numbered line
    (``4.2 Retention``) sets its depth from the numbering, and a short line with no
    terminal punctuation followed by a blank line is treated as a level-2 heading.

    Args:
        builder: Block builder.
        text: Text extracted from the page.
        page: 1-based page number.
    """
    paragraphs = re.split(r"\n\s*\n", text)
    for paragraph in paragraphs:
        lines = [line.strip() for line in paragraph.split("\n") if line.strip()]
        if not lines:
            continue
        first = lines[0]
        numbered = _NUMBERED_HEADING_RE.match(first)
        if numbered and len(lines) == 1:
            depth = numbered.group(1).count(".") + 1
            builder.heading(numbered.group(2), depth, page=page)
            continue
        if (
            len(lines) == 1
            and len(first) <= _HEADING_MAX_CHARS
            and not first.endswith((".", ",", ";", ":"))
            and first[:1].isupper()
        ):
            builder.heading(first, 2, page=page)
            continue
        list_items = [_LIST_RE.match(line) for line in lines]
        if all(match is not None for match in list_items):
            for match in list_items:
                if match is not None:
                    builder.block(BlockKind.LIST_ITEM, match.group(1), page=page)
            continue
        builder.block(BlockKind.PARAGRAPH, " ".join(lines), page=page)


# -------------------------------------------------------------------------- docx
def parse_docx(payload: bytes) -> tuple[_BlockBuilder, dict[str, Any]]:
    """Parse a DOCX file, mapping ``Heading N`` styles onto heading levels.

    Args:
        payload: Raw DOCX bytes.

    Returns:
        A ``(builder, extra)`` pair; ``extra`` may carry ``title`` and ``author``.

    Raises:
        ParseError: If ``python-docx`` is not installed.
    """
    try:
        import docx
    except ImportError as exc:  # pragma: no cover - optional install
        msg = "DOCX parsing needs python-docx"
        raise ParseError(msg) from exc

    document = docx.Document(io.BytesIO(payload))
    builder = _BlockBuilder()
    extra: dict[str, Any] = {}
    properties = document.core_properties
    if properties.title:
        extra["title"] = normalise_text(str(properties.title))
    if properties.author:
        extra["author"] = normalise_text(str(properties.author))

    for paragraph in document.paragraphs:
        style = (paragraph.style.name if paragraph.style is not None else "") or ""
        text = paragraph.text
        if not text.strip():
            continue
        heading_level = _docx_heading_level(style)
        if heading_level is not None:
            builder.heading(text, heading_level)
        elif "List" in style:
            builder.block(BlockKind.LIST_ITEM, text)
        elif "Quote" in style:
            builder.block(BlockKind.QUOTE, text)
        elif "Caption" in style:
            builder.block(BlockKind.CAPTION, text)
        else:
            builder.block(BlockKind.PARAGRAPH, text)

    for table in document.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        builder.block(BlockKind.TABLE, render_table(rows), metadata={"rows": len(rows)})
    return builder, extra


def _docx_heading_level(style: str) -> int | None:
    """Map a DOCX paragraph style onto a heading depth.

    Args:
        style: Style name such as ``"Heading 2"`` or ``"Title"``.

    Returns:
        The heading depth, or None when the style is not a heading.
    """
    lowered = style.lower()
    if lowered in {"title", "subtitle"}:
        return 1
    match = re.match(r"heading\s*(\d)", lowered)
    return int(match.group(1)) if match else None


# -------------------------------------------------------------------------- pptx
def parse_pptx(payload: bytes) -> tuple[_BlockBuilder, dict[str, Any]]:
    """Parse a PPTX deck: one slide becomes one section.

    Args:
        payload: Raw PPTX bytes.

    Returns:
        A ``(builder, extra)`` pair; ``extra`` carries ``page_count`` (slide count).

    Raises:
        ParseError: If ``python-pptx`` is not installed.
    """
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - optional install
        msg = "PPTX parsing needs python-pptx"
        raise ParseError(msg) from exc

    presentation = Presentation(io.BytesIO(payload))
    builder = _BlockBuilder()
    slides = list(presentation.slides)
    for index, slide in enumerate(slides, start=1):
        title = ""
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text
        builder.sections.reset()
        builder.heading(title or f"Slide {index}", 1, page=index)
        for shape in slide.shapes:
            if shape is slide.shapes.title:
                continue
            if getattr(shape, "has_table", False):
                table = shape.table
                rows = [[cell.text for cell in row.cells] for row in table.rows]
                builder.block(
                    BlockKind.TABLE,
                    render_table(rows),
                    page=index,
                    metadata={"rows": len(rows)},
                )
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text
                if not text.strip():
                    continue
                kind = (
                    BlockKind.LIST_ITEM
                    if getattr(paragraph, "level", 0)
                    else BlockKind.PARAGRAPH
                )
                builder.block(kind, text, page=index)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            builder.block(
                BlockKind.FOOTNOTE, slide.notes_slide.notes_text_frame.text, page=index
            )
    return builder, {"page_count": len(slides)}


# -------------------------------------------------------------------------- xlsx
def parse_xlsx(
    payload: bytes, settings: Settings
) -> tuple[_BlockBuilder, dict[str, Any]]:
    """Parse an XLSX workbook: one sheet becomes one section of table blocks.

    Args:
        payload: Raw XLSX bytes.
        settings: Process settings; ``chunk_max_tokens`` bounds one table block.

    Returns:
        A ``(builder, extra)`` pair; ``extra`` carries ``page_count`` (sheet count).

    Raises:
        ParseError: If ``openpyxl`` is not installed.
    """
    try:
        import openpyxl
    except ImportError as exc:  # pragma: no cover - optional install
        msg = "XLSX parsing needs openpyxl"
        raise ParseError(msg) from exc

    workbook = openpyxl.load_workbook(
        io.BytesIO(payload), read_only=True, data_only=True
    )
    builder = _BlockBuilder()
    try:
        for sheet in workbook.worksheets:
            rows = [
                list(row)
                for row in sheet.iter_rows(values_only=True)
                if any(cell is not None and str(cell).strip() for cell in row)
            ]
            if not rows:
                continue
            builder.sections.reset()
            builder.heading(str(sheet.title), 1)
            header, *body = rows
            for batch in _batch_rows(body, header, settings):
                builder.block(
                    BlockKind.TABLE,
                    render_table([header, *batch]),
                    metadata={"sheet": str(sheet.title), "rows": len(batch)},
                )
            if not body:
                builder.block(
                    BlockKind.TABLE,
                    render_table([header]),
                    metadata={"sheet": str(sheet.title), "rows": 0},
                )
        sheet_count = len(workbook.worksheets)
    finally:
        workbook.close()
    return builder, {"page_count": sheet_count}
